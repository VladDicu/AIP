from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import os
import io
import json
from google import genai
from google.genai import types
from pypdf import PdfReader
import docx
import pptx
import openpyxl

app = FastAPI(title="AIP Backend Core")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

os.environ["GEMINI_API_KEY"] = "" # Lasă gol pentru GitHub, vom seta cheia în Render!
client = genai.Client()

# Sistemul de Fallback Automat
def generate_with_fallback(contents, config):
    models_to_try = ["gemini-1.5-flash", "gemini-1.5-pro", "gemini-2.0-flash"]
    
    for model_id in models_to_try:
        try:
            return client.models.generate_content(
                model=model_id,
                contents=contents,
                config=config
            )
        except Exception as e:
            if "503" in str(e) or "UNAVAILABLE" in str(e):
                continue
            raise e
            
    raise Exception("Toate serverele AI sunt momentan suprasolicitate. Reîncearcă în câteva secunde.")

class QueryRequest(BaseModel):
    query: str
    strict_mode: bool
    context: str = ""

class ActionRequest(BaseModel):
    context: str
    action_type: str

@app.post("/api/upload")
async def upload_document(file: UploadFile = File(...)):
    try:
        content = await file.read()
        text = ""
        filename = file.filename.lower()
        
        if filename.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(content))
            text = "".join([page.extract_text() or "" for page in reader.pages])
            
        elif filename.endswith((".docx", ".doc")):
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs])
            
        elif filename.endswith(".pptx"):
            prs = pptx.Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif filename.endswith((".xlsx", ".xls", ".csv")):
            if filename.endswith(".csv"):
                text = content.decode('utf-8', errors='ignore')
            else:
                wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = [str(c) for c in row if c is not None]
                        if row_text:
                            text += " | ".join(row_text) + "\n"
        else:
            try:
                text = content.decode('utf-8')
            except UnicodeDecodeError:
                text = content.decode('latin-1', errors='ignore')

        if not text.strip():
            return {"error": "Fișierul pare gol sau conține doar elemente grafice fără text."}
            
        return {
            "filename": file.filename,
            "status": "Document procesat cu succes",
            "extracted_text": text 
        }
    except Exception as e:
        return {"error": f"Eroare la procesare: {str(e)}"}

@app.post("/api/ask")
def process_query(request: QueryRequest):
    if request.strict_mode:
        sys_instruct = "Ești un asistent academic. Răspunde STRICT pe baza contextului oferit."
        prompt = f"Context:\n{request.context}\n\nÎntrebare: {request.query}"
        temp = 0.1
        confidence_msg = "100% (Strict Source)"
    else:
        sys_instruct = "Ești un asistent academic. Răspunde folosind și baza ta de cunoștințe. Adaugă un Confidence Score la final."
        prompt = f"Întrebare: {request.query}"
        temp = 0.7
        confidence_msg = "Evaluat de AI in text"

    try:
        response = generate_with_fallback(
            contents=prompt,
            config=types.GenerateContentConfig(system_instruction=sys_instruct, temperature=temp)
        )
        return {"answer": response.text, "strict_mode_active": request.strict_mode, "confidence_status": confidence_msg}
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/action")
def fast_action(request: ActionRequest):
    if not request.context:
        return {"error": "Încarcă un document mai întâi."}
        
    if request.action_type == "summary":
        prompt = f"Extrage un rezumat clar, în 3-5 idei principale, din:\n{request.context}"
        try:
            response = generate_with_fallback(contents=prompt, config=types.GenerateContentConfig(temperature=0.3))
            return {"answer": response.text}
        except Exception as e:
            return {"error": str(e)}
            
    elif request.action_type == "quiz":
        prompt = f"""Creează un test interactiv cu 3 întrebări din următorul curs.
        Răspunde STRICT în format JSON valid:
        [
          {{
            "question": "Text?",
            "options": ["V1", "V2", "V3", "V4"],
            "correct_index": 0
          }}
        ]
        Curs:
        {request.context}"""
        
        try:
            response = generate_with_fallback(contents=prompt, config=types.GenerateContentConfig(temperature=0.2))
            raw_text = response.text.strip()
            if raw_text.startswith("```json"): raw_text = raw_text[7:-3].strip()
            elif raw_text.startswith("```"): raw_text = raw_text[3:-3].strip()
                
            return {"type": "quiz", "data": json.loads(raw_text)}
        except Exception as e:
            return {"error": str(e)}